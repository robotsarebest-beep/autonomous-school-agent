import docx
import pptx
import pdfplumber
import os
from core.logger import logger

class FileProcessor:
    @staticmethod
    def extract_docx(file_path):
        try:
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return ""

    @staticmethod
    def extract_pptx(file_path):
        try:
            prs = pptx.Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return ""

    @staticmethod
    def extract_pdf(file_path):
        try:
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text.append(page.extract_text() or "")
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return ""

    def process_any(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".docx":
            return self.extract_docx(file_path)
        elif ext == ".pptx":
            return self.extract_pptx(file_path)
        elif ext == ".pdf":
            return self.extract_pdf(file_path)
        else:
            logger.warn(f"Unsupported file format: {ext}")
            return ""
